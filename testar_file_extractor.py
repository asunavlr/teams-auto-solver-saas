"""
Testes para engine/file_extractor.py
Cria arquivos temporarios de cada tipo e verifica a extracao.
"""

import sys
import tempfile
from pathlib import Path

# Adiciona o diretorio raiz ao path
sys.path.insert(0, str(Path(__file__).parent))

from engine.file_extractor import (
    extrair_conteudo_pdf,
    extrair_conteudo_docx,
    extrair_conteudo_xlsx,
    extrair_conteudo_pptx,
    extrair_conteudo_arquivo,
)

PASSED = 0
FAILED = 0


def test(nome, condicao, detalhe=""):
    global PASSED, FAILED
    if condicao:
        PASSED += 1
        print(f"  [OK] {nome}")
    else:
        FAILED += 1
        print(f"  [FALHOU] {nome} - {detalhe}")


def testar_docx():
    """Testa extracao de DOCX."""
    print("\n=== Teste DOCX ===")
    from docx import Document

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        doc = Document()
        doc.add_heading("Titulo do Trabalho", 0)
        doc.add_paragraph("Este e o primeiro paragrafo do documento.")
        doc.add_paragraph("Segundo paragrafo com mais conteudo sobre o tema.")

        # Adiciona tabela
        table = doc.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Nome"
        table.cell(0, 1).text = "Nota"
        table.cell(1, 0).text = "Joao"
        table.cell(1, 1).text = "8.5"
        table.cell(2, 0).text = "Maria"
        table.cell(2, 1).text = "9.0"

        doc.save(tmp.name)
        filepath = Path(tmp.name)

    resultado = extrair_conteudo_docx(filepath)
    test("Retorna dict", isinstance(resultado, dict))
    test("Tem chave 'texto'", "texto" in resultado)
    test("Tem chave 'paginas'", "paginas" in resultado)
    test("Texto nao vazio", len(resultado["texto"]) > 0, f"len={len(resultado['texto'])}")
    test("Contem titulo", "Titulo do Trabalho" in resultado["texto"], resultado["texto"][:100])
    test("Contem paragrafo", "primeiro paragrafo" in resultado["texto"])
    test("Contem dados da tabela", "Joao" in resultado["texto"] and "8.5" in resultado["texto"])

    filepath.unlink()


def testar_xlsx():
    """Testa extracao de XLSX."""
    print("\n=== Teste XLSX ===")
    from openpyxl import Workbook

    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
        wb = Workbook()
        ws = wb.active
        ws.title = "Notas"
        ws.append(["Aluno", "Prova 1", "Prova 2", "Media"])
        ws.append(["Joao Silva", 7.5, 8.0, 7.75])
        ws.append(["Maria Santos", 9.0, 8.5, 8.75])
        ws.append(["Pedro Costa", 6.0, 7.0, 6.5])
        wb.save(tmp.name)
        filepath = Path(tmp.name)

    resultado = extrair_conteudo_xlsx(filepath)
    test("Retorna dict", isinstance(resultado, dict))
    test("Texto nao vazio", len(resultado["texto"]) > 0)
    test("Contem nome planilha", "Notas" in resultado["texto"])
    test("Contem dados", "Joao Silva" in resultado["texto"])
    test("Contem numeros", "7.5" in resultado["texto"])

    filepath.unlink()


def testar_pptx():
    """Testa extracao de PPTX."""
    print("\n=== Teste PPTX ===")
    from pptx import Presentation
    from pptx.util import Inches

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs = Presentation()

        # Slide 1 - titulo
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Apresentacao de Teste"
        slide.placeholders[1].text = "Subtitulo da apresentacao"

        # Slide 2 - conteudo
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Slide de Conteudo"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = "Primeiro ponto importante"
        p = tf.add_paragraph()
        p.text = "Segundo ponto sobre o tema"

        prs.save(tmp.name)
        filepath = Path(tmp.name)

    resultado = extrair_conteudo_pptx(filepath)
    test("Retorna dict", isinstance(resultado, dict))
    test("Texto nao vazio", len(resultado["texto"]) > 0)
    test("Paginas = 2", resultado["paginas"] == 2, f"paginas={resultado['paginas']}")
    test("Contem titulo", "Apresentacao de Teste" in resultado["texto"])
    test("Contem conteudo slide 2", "ponto importante" in resultado["texto"])
    test("Tem separador de slides", "Slide 1" in resultado["texto"] and "Slide 2" in resultado["texto"])

    filepath.unlink()


def testar_pdf():
    """Testa extracao de PDF."""
    print("\n=== Teste PDF ===")
    import pdfplumber

    # Cria um PDF simples usando pdfplumber/pdfminer nao cria PDFs,
    # entao vamos testar com um PDF gerado via reportlab ou fpdf se disponivel
    # Fallback: testa com um arquivo que nao existe (deve retornar vazio sem crash)

    # Testa resiliencia com arquivo inexistente
    resultado_fake = extrair_conteudo_pdf(Path("/tmp/nao_existe.pdf"))
    test("PDF inexistente nao crasha", isinstance(resultado_fake, dict))
    test("PDF inexistente retorna texto vazio", resultado_fake["texto"] == "")
    test("PDF inexistente retorna base64 None", resultado_fake["base64_data"] is None)

    # Tenta criar PDF real se fpdf2 estiver disponivel
    try:
        from fpdf import FPDF

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Helvetica", size=12)
            pdf.cell(200, 10, text="Titulo do Documento PDF", new_x="LMARGIN", new_y="NEXT")
            pdf.cell(200, 10, text="Conteudo da primeira pagina com informacoes.", new_x="LMARGIN", new_y="NEXT")
            pdf.cell(200, 10, text="Mais texto para testar extracao.", new_x="LMARGIN", new_y="NEXT")

            pdf.add_page()
            pdf.cell(200, 10, text="Segunda pagina do documento.", new_x="LMARGIN", new_y="NEXT")
            pdf.cell(200, 10, text="Exercicio 1: Responda a questao.", new_x="LMARGIN", new_y="NEXT")

            pdf.output(tmp.name)
            filepath = Path(tmp.name)

        resultado = extrair_conteudo_pdf(filepath)
        test("PDF real: texto nao vazio", len(resultado["texto"]) > 0, f"len={len(resultado['texto'])}")
        test("PDF real: 2 paginas", resultado["paginas"] == 2, f"paginas={resultado['paginas']}")
        test("PDF real: base64 gerado", resultado["base64_data"] is not None)
        test("PDF real: base64 nao vazio", len(resultado["base64_data"]) > 0 if resultado["base64_data"] else False)
        test("PDF real: contem titulo", "Titulo" in resultado["texto"], resultado["texto"][:200])
        test("PDF real: contem segunda pagina", "Segunda pagina" in resultado["texto"] or "Pagina 2" in resultado["texto"])

        filepath.unlink()

    except ImportError:
        print("  [SKIP] fpdf2 nao instalado, pulando teste de PDF real")
        print("         Instale com: pip install fpdf2")


def testar_dispatcher():
    """Testa o dispatcher extrair_conteudo_arquivo."""
    print("\n=== Teste Dispatcher ===")
    from docx import Document

    # Testa com DOCX
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        doc = Document()
        doc.add_paragraph("Teste do dispatcher")
        doc.save(tmp.name)
        filepath = Path(tmp.name)

    resultado = extrair_conteudo_arquivo(filepath)
    test("Dispatcher DOCX retorna dict", isinstance(resultado, dict))
    test("Dispatcher DOCX contem texto", "dispatcher" in resultado.get("texto", "").lower())
    filepath.unlink()

    # Testa com extensao nao suportada
    with tempfile.NamedTemporaryFile(suffix=".xyz", delete=False) as tmp:
        tmp.write(b"conteudo qualquer")
        filepath = Path(tmp.name)

    resultado = extrair_conteudo_arquivo(filepath)
    test("Dispatcher extensao desconhecida retorna None", resultado is None)
    filepath.unlink()


def testar_imports_monitor():
    """Testa se os imports do monitor.py funcionam."""
    print("\n=== Teste Imports ===")
    try:
        from engine.file_extractor import extrair_conteudo_arquivo
        test("Import file_extractor OK", True)
    except Exception as e:
        test("Import file_extractor OK", False, str(e))

    try:
        from engine.solver import resolver_com_claude
        test("Import solver OK", True)
    except Exception as e:
        test("Import solver OK", False, str(e))

    try:
        from engine.agent import TeamsAgent
        test("Import agent OK", True)
        test("Seletor download_preview existe", "download_preview" in TeamsAgent.SELECTORS)
        test("Seletor menu_tres_pontos existe", "menu_tres_pontos" in TeamsAgent.SELECTORS)
        test("Seletor download_menu_item existe", "download_menu_item" in TeamsAgent.SELECTORS)
        test("Descricao download_preview existe", "download_preview" in TeamsAgent.DESCRICOES)
        test("Descricao menu_tres_pontos existe", "menu_tres_pontos" in TeamsAgent.DESCRICOES)
        test("Descricao download_menu_item existe", "download_menu_item" in TeamsAgent.DESCRICOES)
    except Exception as e:
        test("Import agent OK", False, str(e))

    try:
        from engine.browser import TeamsBrowser
        test("Import browser OK", True)
    except Exception as e:
        test("Import browser OK", False, str(e))


def testar_solver_aceita_novos_campos():
    """Testa se o solver aceita os novos campos do tarefa_info."""
    print("\n=== Teste Solver (novos campos) ===")
    import inspect
    from engine.solver import resolver_com_claude

    # Verifica se o codigo do solver referencia os novos campos
    source = inspect.getsource(resolver_com_claude)
    test("Solver referencia texto_extraido", "texto_extraido" in source)
    test("Solver referencia pdf_base64", "pdf_base64" in source)
    test("Solver usa type document", '"document"' in source or "'document'" in source)


if __name__ == "__main__":
    print("=" * 60)
    print("TESTES - file_extractor.py + integracoes")
    print("=" * 60)

    testar_imports_monitor()
    testar_docx()
    testar_xlsx()
    testar_pptx()
    testar_pdf()
    testar_dispatcher()
    testar_solver_aceita_novos_campos()

    print("\n" + "=" * 60)
    print(f"RESULTADO: {PASSED} passed, {FAILED} failed")
    print("=" * 60)

    sys.exit(1 if FAILED > 0 else 0)

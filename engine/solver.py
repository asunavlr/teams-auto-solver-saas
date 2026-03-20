"""
Modulo de resolucao de tarefas com Claude API.
Inclui deteccao de formato e criacao de arquivos.
"""

import base64
import json
import os
import re
import time
import zipfile
from datetime import datetime
from pathlib import Path

import anthropic
from docx import Document
from docx.shared import Pt
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
from pptx import Presentation
from pptx.util import Pt as PptxPt
from loguru import logger

CLAUDE_MAX_RETRIES = 3
FORMATOS_CODIGO = ["py", "js", "ts", "java", "c", "cpp", "css", "sql", "kotlin"]
FORMATOS_ARQUIVO = ["html", "docx", "xlsx", "pptx", "zip", "android"] + FORMATOS_CODIGO
TODOS_FORMATOS = FORMATOS_ARQUIVO + ["txt", "texto"]


def detectar_formato_resposta(content: str) -> str:
    """Detecta o formato esperado baseado nas instrucoes da pagina."""
    content_lower = content.lower()

    if any(x in content_lower for x in [".html", "arquivo html", "html file", "pagina html", "codigo html"]):
        return "html"
    if any(x in content_lower for x in [".docx", ".doc", "word", "documento word", "arquivo word"]):
        return "docx"
    if any(x in content_lower for x in [".xlsx", ".xls", "excel", "planilha", "spreadsheet"]):
        return "xlsx"
    if any(x in content_lower for x in [".pptx", ".ppt", "powerpoint", "apresentacao", "slides"]):
        return "pptx"
    if any(x in content_lower for x in [".py", "python", "codigo python", "script python"]):
        return "py"
    if any(x in content_lower for x in [".js", "javascript"]):
        return "js"
    if any(x in content_lower for x in [".ts", "typescript"]):
        return "ts"
    if any(x in content_lower for x in [".java", "codigo java", "programa java"]):
        return "java"
    if any(x in content_lower for x in ["linguagem c ", "codigo em c ", "programa em c ", "arquivo .c"]):
        return "c"
    if any(x in content_lower for x in [".cpp", "c++", "codigo c++", "programa c++"]):
        return "cpp"
    if any(x in content_lower for x in [".css", "arquivo css", "folha de estilo"]):
        return "css"
    if any(x in content_lower for x in [".sql", "script sql", "consulta sql", "query sql"]):
        return "sql"
    if any(x in content_lower for x in [".kt", ".kotlin", "kotlin"]):
        return "kotlin"
    if any(x in content_lower for x in [
        "android studio", "android app", "aplicativo android", "app android",
        "activity", "layout xml", "androidmanifest", "projeto android"
    ]):
        return "android"
    if any(x in content_lower for x in [".zip", "arquivo zip", "compactado"]):
        return "zip"

    return "texto"


def detectar_formato_da_resposta(resposta: str) -> str | None:
    """Detecta o formato da resposta do Claude."""
    resposta_lower = resposta.lower()

    formatos_validos = "|".join(TODOS_FORMATOS)
    match = re.search(rf'\[formato:\s*({formatos_validos})\]', resposta_lower)
    if match:
        return match.group(1)

    if any(tag in resposta_lower for tag in ["<!doctype html", "<html", "<head>", "<body>"]):
        return "html"
    if "```html" in resposta_lower:
        return "html"
    if "```python" in resposta_lower or "```py" in resposta_lower:
        return "py"
    if "```javascript" in resposta_lower or "```js" in resposta_lower:
        return "js"
    if "```typescript" in resposta_lower or "```ts" in resposta_lower:
        return "ts"
    if "```java" in resposta_lower:
        return "java"
    if "```cpp" in resposta_lower or "```c++" in resposta_lower:
        return "cpp"
    if "```c\n" in resposta_lower:
        return "c"
    if "```css" in resposta_lower:
        return "css"
    if "```sql" in resposta_lower:
        return "sql"
    if "```kotlin" in resposta_lower or "```kt" in resposta_lower:
        return "kotlin"
    if "```xml" in resposta_lower and ("android" in resposta_lower or "layout" in resposta_lower):
        return "android"

    return None


def remover_marcador_formato(resposta: str) -> str:
    return re.sub(r'\[FORMATO:\s*\w+\]\s*\n?', '', resposta, count=1, flags=re.IGNORECASE).strip()


# ============================================================
# CRIACAO DE ARQUIVOS
# ============================================================

def criar_arquivo_html(conteudo: str, nome_arquivo: str, data_dir: Path) -> str:
    if "<html" in conteudo.lower() or "<!doctype" in conteudo.lower():
        html_content = conteudo
    else:
        html_content = f"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{nome_arquivo}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 20px; line-height: 1.6; }}
        h1, h2, h3 {{ color: #333; }}
    </style>
</head>
<body>
{conteudo}
</body>
</html>"""

    filepath = data_dir / f"{nome_arquivo}.html"
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(html_content)
    return str(filepath)


def criar_arquivo_docx(conteudo: str, nome_arquivo: str, data_dir: Path) -> str:
    doc = Document()
    doc.add_heading(nome_arquivo, 0)

    for linha in conteudo.split('\n'):
        linha_stripped = linha.strip()
        if not linha_stripped:
            doc.add_paragraph("")
            continue
        if linha_stripped.startswith('# '):
            doc.add_heading(linha_stripped[2:], level=1)
        elif linha_stripped.startswith('## '):
            doc.add_heading(linha_stripped[3:], level=2)
        elif linha_stripped.startswith('### '):
            doc.add_heading(linha_stripped[4:], level=3)
        elif linha_stripped.startswith('- ') or linha_stripped.startswith('* '):
            doc.add_paragraph(linha_stripped[2:], style='List Bullet')
        elif re.match(r'^\d+\.\s', linha_stripped):
            texto = re.sub(r'^\d+\.\s', '', linha_stripped)
            doc.add_paragraph(texto, style='List Number')
        else:
            p = doc.add_paragraph(linha_stripped)
            p.style.font.size = Pt(11)

    filepath = data_dir / f"{nome_arquivo}.docx"
    doc.save(str(filepath))
    return str(filepath)


def criar_arquivo_xlsx(conteudo: str, nome_arquivo: str, data_dir: Path) -> str:
    wb = Workbook()
    ws = wb.active
    ws.title = nome_arquivo[:31]

    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font_white = Font(bold=True, size=11, color="FFFFFF")
    thin_border = Border(
        left=Side(style='thin'), right=Side(style='thin'),
        top=Side(style='thin'), bottom=Side(style='thin')
    )

    linhas = conteudo.strip().split('\n')
    is_first_data_row = True

    for linha in linhas:
        linha = linha.strip()
        if not linha or re.match(r'^[\|\-\s:]+$', linha):
            continue

        if '|' in linha:
            cells = [c.strip() for c in linha.split('|') if c.strip()]
            for j, cell in enumerate(cells):
                c = ws.cell(row=ws.max_row + (0 if is_first_data_row else 1), column=j + 1, value=cell)
                c.border = thin_border
                if is_first_data_row:
                    c.font = header_font_white
                    c.fill = header_fill
                    c.alignment = Alignment(horizontal='center')
            is_first_data_row = False
        else:
            ws.cell(row=ws.max_row + 1, column=1, value=linha)

    for col in ws.columns:
        max_length = 0
        col_letter = col[0].column_letter
        for cell in col:
            if cell.value:
                max_length = max(max_length, len(str(cell.value)))
        ws.column_dimensions[col_letter].width = min(max_length + 4, 50)

    filepath = data_dir / f"{nome_arquivo}.xlsx"
    wb.save(str(filepath))
    return str(filepath)


def criar_arquivo_pptx(conteudo: str, nome_arquivo: str, data_dir: Path) -> str:
    prs = Presentation()

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = nome_arquivo
    slide.placeholders[1].text = datetime.now().strftime("%d/%m/%Y")

    slides_content = []
    current_slide = {"titulo": "", "conteudo": []}

    for linha in conteudo.split('\n'):
        linha = linha.strip()
        if linha.startswith('# ') or linha.startswith('## '):
            if current_slide["titulo"] or current_slide["conteudo"]:
                slides_content.append(current_slide)
            titulo = re.sub(r'^#+\s*', '', linha)
            current_slide = {"titulo": titulo, "conteudo": []}
        elif linha:
            current_slide["conteudo"].append(linha)

    if current_slide["titulo"] or current_slide["conteudo"]:
        slides_content.append(current_slide)

    for sc in slides_content:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = sc["titulo"]
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, linha in enumerate(sc["conteudo"]):
            if i == 0:
                tf.text = linha
            else:
                p = tf.add_paragraph()
                p.text = linha
                p.font.size = PptxPt(16)

    filepath = data_dir / f"{nome_arquivo}.pptx"
    prs.save(str(filepath))
    return str(filepath)


def criar_arquivo_codigo(conteudo: str, nome_arquivo: str, extensao: str, data_dir: Path) -> str:
    conteudo = re.sub(r'^```\w*\n?', '', conteudo)
    conteudo = re.sub(r'\n?```$', '', conteudo)

    filepath = data_dir / f"{nome_arquivo}.{extensao}"
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(conteudo)
    return str(filepath)


def criar_arquivo_zip(arquivos: list, nome_arquivo: str, data_dir: Path) -> str:
    filepath = data_dir / f"{nome_arquivo}.zip"
    with zipfile.ZipFile(filepath, 'w', zipfile.ZIP_DEFLATED) as zf:
        for arq in arquivos:
            zf.write(arq, os.path.basename(arq))
    return str(filepath)


def criar_arquivo_resposta(conteudo: str, nome_tarefa: str, formato: str, data_dir: Path) -> str:
    nome_limpo = re.sub(r'[^\w\s-]', '', nome_tarefa)[:50].strip().replace(' ', '_')

    if formato == "html":
        return criar_arquivo_html(conteudo, nome_limpo, data_dir)
    elif formato == "docx":
        return criar_arquivo_docx(conteudo, nome_limpo, data_dir)
    elif formato == "xlsx":
        return criar_arquivo_xlsx(conteudo, nome_limpo, data_dir)
    elif formato == "pptx":
        return criar_arquivo_pptx(conteudo, nome_limpo, data_dir)
    elif formato in FORMATOS_CODIGO:
        return criar_arquivo_codigo(conteudo, nome_limpo, formato, data_dir)
    else:
        filepath = data_dir / f"{nome_limpo}.txt"
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(conteudo)
        return str(filepath)


def extrair_multiplos_arquivos(resposta: str, formato: str, nome_tarefa: str, data_dir: Path) -> list:
    """Extrai multiplos arquivos da resposta do Claude."""
    arquivos = []
    nome_limpo = re.sub(r'[^\w\s-]', '', nome_tarefa)[:50].strip().replace(' ', '_')

    if formato == "html":
        blocos_html = re.findall(r'```html\s*\n(.*?)```', resposta, re.DOTALL)
        if not blocos_html:
            blocos_html = re.findall(r'(<!DOCTYPE html>.*?</html>)', resposta, re.DOTALL | re.IGNORECASE)
        if not blocos_html:
            blocos_html = [resposta]

        for i, bloco in enumerate(blocos_html):
            bloco = bloco.strip()
            if not bloco:
                continue
            if "<html" not in bloco.lower():
                bloco = f"""<!DOCTYPE html>
<html lang="pt-BR">
<head><meta charset="UTF-8"><title>{nome_tarefa} - Exercicio {i+1}</title></head>
<body>{bloco}</body>
</html>"""
            sufixo = f"_{i+1}" if len(blocos_html) > 1 else ""
            filepath = data_dir / f"{nome_limpo}{sufixo}.html"
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(bloco)
            arquivos.append(str(filepath))

    elif formato in FORMATOS_CODIGO:
        lang_map = {
            "py": ["python", "py"], "js": ["javascript", "js"],
            "ts": ["typescript", "ts"], "java": ["java"], "c": ["c"],
            "cpp": ["cpp", "c\\+\\+"], "css": ["css"], "sql": ["sql"],
        }
        langs = lang_map.get(formato, [formato])
        pattern = r'```(?:' + '|'.join(langs) + r')\s*\n(.*?)```'
        blocos = re.findall(pattern, resposta, re.DOTALL)

        if not blocos:
            cleaned = re.sub(r'^```\w*\n?', '', resposta)
            cleaned = re.sub(r'\n?```$', '', cleaned)
            blocos = [cleaned]

        for i, bloco in enumerate(blocos):
            bloco = bloco.strip()
            if not bloco:
                continue
            sufixo = f"_{i+1}" if len(blocos) > 1 else ""
            filepath = data_dir / f"{nome_limpo}{sufixo}.{formato}"
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(bloco)
            arquivos.append(str(filepath))

    return arquivos


def extrair_projeto_multi_arquivo(resposta: str, nome_tarefa: str, data_dir: Path) -> list:
    """Extrai projeto multi-arquivo e zipa."""
    arquivos = []
    nome_limpo = re.sub(r'[^\w\s-]', '', nome_tarefa)[:50].strip().replace(' ', '_')

    blocos = re.findall(r'```(\w+)\s*\n(.*?)```', resposta, re.DOTALL)
    ext_map = {
        "html": "html", "css": "css", "javascript": "js", "js": "js",
        "typescript": "ts", "ts": "ts", "python": "py", "py": "py",
        "java": "java", "c": "c", "cpp": "cpp", "sql": "sql",
        "kotlin": "kt", "kt": "kt", "xml": "xml",
    }

    contadores = {}
    for lang, conteudo in blocos:
        ext = ext_map.get(lang.lower(), lang.lower())
        conteudo = conteudo.strip()
        if not conteudo:
            continue
        contadores[ext] = contadores.get(ext, 0) + 1
        sufixo = f"_{contadores[ext]}" if contadores[ext] > 1 else ""
        filepath = data_dir / f"{nome_limpo}{sufixo}.{ext}"
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(conteudo)
        arquivos.append(str(filepath))

    # Cria ZIP se houver mais de 1 arquivo (independente do tipo)
    if len(arquivos) > 1:
        zip_path = criar_arquivo_zip(arquivos, nome_limpo, data_dir)
        return [zip_path]

    return arquivos


def criar_projeto_android(resposta: str, nome_tarefa: str, data_dir: Path) -> list:
    """
    Cria projeto Android Studio completo a partir da resposta do Claude.

    Espera que Claude gere blocos marcados com:
    - ```java ou ```kotlin para código
    - ```xml para layouts e manifest

    Estrutura gerada:
    NomeProjeto/
    ├── app/src/main/java/com/example/app/
    ├── app/src/main/res/layout/
    ├── app/src/main/AndroidManifest.xml
    └── outros arquivos de configuração
    """
    import shutil

    # Nome do projeto limpo
    nome_limpo = re.sub(r'[^\w\s-]', '', nome_tarefa)[:30].strip().replace(' ', '_')
    nome_projeto = nome_limpo or "MeuApp"
    package_name = "com.example." + nome_projeto.lower().replace('_', '').replace('-', '')
    package_path = package_name.replace('.', '/')

    # Diretório temporário do projeto
    projeto_dir = data_dir / nome_projeto
    if projeto_dir.exists():
        shutil.rmtree(projeto_dir)

    # Estrutura de diretórios
    dirs = [
        projeto_dir / "app" / "src" / "main" / "java" / package_path,
        projeto_dir / "app" / "src" / "main" / "res" / "layout",
        projeto_dir / "app" / "src" / "main" / "res" / "values",
        projeto_dir / "app" / "src" / "main" / "res" / "drawable",
    ]
    for d in dirs:
        d.mkdir(parents=True, exist_ok=True)

    # Parse dos blocos de código
    blocos = re.findall(r'```(\w+)\s*\n(.*?)```', resposta, re.DOTALL)

    java_files = []
    kotlin_files = []
    xml_files = []

    for lang, conteudo in blocos:
        conteudo = conteudo.strip()
        if not conteudo:
            continue
        lang_lower = lang.lower()

        if lang_lower == "java":
            java_files.append(conteudo)
        elif lang_lower in ["kotlin", "kt"]:
            kotlin_files.append(conteudo)
        elif lang_lower == "xml":
            xml_files.append(conteudo)

    # Detecta nome da classe principal nos arquivos Java/Kotlin
    def extrair_nome_classe(codigo):
        match = re.search(r'class\s+(\w+)', codigo)
        return match.group(1) if match else None

    # Salva arquivos Java
    for i, codigo in enumerate(java_files):
        nome_classe = extrair_nome_classe(codigo) or f"Classe{i+1}"
        filepath = projeto_dir / "app" / "src" / "main" / "java" / package_path / f"{nome_classe}.java"
        # Adiciona package se não tiver
        if "package " not in codigo:
            codigo = f"package {package_name};\n\n{codigo}"
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(codigo)

    # Salva arquivos Kotlin
    for i, codigo in enumerate(kotlin_files):
        nome_classe = extrair_nome_classe(codigo) or f"Classe{i+1}"
        filepath = projeto_dir / "app" / "src" / "main" / "java" / package_path / f"{nome_classe}.kt"
        # Adiciona package se não tiver
        if "package " not in codigo:
            codigo = f"package {package_name}\n\n{codigo}"
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(codigo)

    # Salva arquivos XML
    for i, codigo in enumerate(xml_files):
        if "AndroidManifest" in codigo or "manifest" in codigo.lower()[:100]:
            filepath = projeto_dir / "app" / "src" / "main" / "AndroidManifest.xml"
        elif "layout" in codigo.lower()[:200] or "LinearLayout" in codigo or "RelativeLayout" in codigo or "ConstraintLayout" in codigo:
            # Tenta extrair nome do layout
            nome_layout = f"activity_main" if i == 0 else f"layout_{i+1}"
            filepath = projeto_dir / "app" / "src" / "main" / "res" / "layout" / f"{nome_layout}.xml"
        elif "resources" in codigo.lower()[:100] and "string" in codigo.lower():
            filepath = projeto_dir / "app" / "src" / "main" / "res" / "values" / "strings.xml"
        else:
            filepath = projeto_dir / "app" / "src" / "main" / "res" / "layout" / f"layout_{i+1}.xml"
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(codigo)

    # Gera AndroidManifest.xml se não foi fornecido
    manifest_path = projeto_dir / "app" / "src" / "main" / "AndroidManifest.xml"
    if not manifest_path.exists():
        # Detecta a Activity principal
        main_activity = "MainActivity"
        for codigo in java_files + kotlin_files:
            if "AppCompatActivity" in codigo or "Activity" in codigo:
                nome = extrair_nome_classe(codigo)
                if nome:
                    main_activity = nome
                    break

        manifest_content = f'''<?xml version="1.0" encoding="utf-8"?>
<manifest xmlns:android="http://schemas.android.com/apk/res/android"
    package="{package_name}">

    <application
        android:allowBackup="true"
        android:icon="@mipmap/ic_launcher"
        android:label="@string/app_name"
        android:roundIcon="@mipmap/ic_launcher_round"
        android:supportsRtl="true"
        android:theme="@style/Theme.AppCompat.Light.DarkActionBar">
        <activity
            android:name=".{main_activity}"
            android:exported="true">
            <intent-filter>
                <action android:name="android.intent.action.MAIN" />
                <category android:name="android.intent.category.LAUNCHER" />
            </intent-filter>
        </activity>
    </application>

</manifest>
'''
        with open(manifest_path, "w", encoding="utf-8") as f:
            f.write(manifest_content)

    # Gera strings.xml se não foi fornecido
    strings_path = projeto_dir / "app" / "src" / "main" / "res" / "values" / "strings.xml"
    if not strings_path.exists():
        strings_content = f'''<?xml version="1.0" encoding="utf-8"?>
<resources>
    <string name="app_name">{nome_projeto}</string>
</resources>
'''
        with open(strings_path, "w", encoding="utf-8") as f:
            f.write(strings_content)

    # Gera build.gradle (app level)
    build_gradle_app = f'''plugins {{
    id 'com.android.application'
}}

android {{
    namespace '{package_name}'
    compileSdk 34

    defaultConfig {{
        applicationId "{package_name}"
        minSdk 24
        targetSdk 34
        versionCode 1
        versionName "1.0"
    }}

    buildTypes {{
        release {{
            minifyEnabled false
        }}
    }}
    compileOptions {{
        sourceCompatibility JavaVersion.VERSION_1_8
        targetCompatibility JavaVersion.VERSION_1_8
    }}
}}

dependencies {{
    implementation 'androidx.appcompat:appcompat:1.6.1'
    implementation 'com.google.android.material:material:1.11.0'
    implementation 'androidx.constraintlayout:constraintlayout:2.1.4'
}}
'''
    with open(projeto_dir / "app" / "build.gradle", "w", encoding="utf-8") as f:
        f.write(build_gradle_app)

    # Gera build.gradle (project level)
    build_gradle_project = '''plugins {
    id 'com.android.application' version '8.2.0' apply false
}
'''
    with open(projeto_dir / "build.gradle", "w", encoding="utf-8") as f:
        f.write(build_gradle_project)

    # Gera settings.gradle
    settings_gradle = f'''pluginManagement {{
    repositories {{
        google()
        mavenCentral()
        gradlePluginPortal()
    }}
}}
dependencyResolutionManagement {{
    repositoriesMode.set(RepositoriesMode.FAIL_ON_PROJECT_REPOS)
    repositories {{
        google()
        mavenCentral()
    }}
}}

rootProject.name = "{nome_projeto}"
include ':app'
'''
    with open(projeto_dir / "settings.gradle", "w", encoding="utf-8") as f:
        f.write(settings_gradle)

    # Gera gradle.properties
    gradle_properties = '''android.useAndroidX=true
android.enableJetifier=true
'''
    with open(projeto_dir / "gradle.properties", "w", encoding="utf-8") as f:
        f.write(gradle_properties)

    # Zipa o projeto
    zip_path = data_dir / f"{nome_projeto}.zip"
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for root, dirs, files in os.walk(projeto_dir):
            for file in files:
                file_path = Path(root) / file
                arcname = file_path.relative_to(data_dir)
                zipf.write(file_path, arcname)

    # Remove diretório temporário
    shutil.rmtree(projeto_dir)

    logger.info(f"Projeto Android criado: {zip_path}")
    return [str(zip_path)]


# ============================================================
# ANALISE DE INTENCAO DA TAREFA
# ============================================================

# Categorias de tarefas
CATEGORIA_RESOLVIVEL = "RESOLVIVEL"
CATEGORIA_RESOLVIVEL_MANUAL = "RESOLVIVEL_MANUAL"  # Resolver mas não enviar
CATEGORIA_RESOLVER_PARCIAL = "RESOLVER_PARCIAL"  # Resolver exercicios, aluno adiciona fotos
CATEGORIA_CERTIFICADO = "CERTIFICADO"
CATEGORIA_AVISO = "AVISO"
CATEGORIA_GRUPO = "GRUPO"
CATEGORIA_RECURSO_EXTERNO = "RECURSO_EXTERNO"
CATEGORIA_PRESENCIAL = "PRESENCIAL"
CATEGORIA_PESSOAL = "PESSOAL"
CATEGORIA_INCERTO = "INCERTO"

CATEGORIAS_PULAR = [
    CATEGORIA_CERTIFICADO,
    CATEGORIA_AVISO,
    CATEGORIA_GRUPO,
    CATEGORIA_RECURSO_EXTERNO,
    CATEGORIA_PRESENCIAL,
    CATEGORIA_PESSOAL,
]

# Categorias que resolvemos mas não enviamos automaticamente
CATEGORIAS_ANEXAR_APENAS = [
    CATEGORIA_RESOLVIVEL_MANUAL,
    CATEGORIA_RESOLVER_PARCIAL,
]


def analisar_intencao_tarefa(tarefa: dict, api_key: str) -> dict:
    """
    Analisa a intencao da tarefa e determina se deve ser resolvida.

    Args:
        tarefa: Dict com nome, instrucoes, screenshots, texto_extraido
        api_key: Chave da API Anthropic

    Returns:
        {
            "categoria": str,  # RESOLVIVEL, AVISO, PESSOAL, etc
            "confianca": int,  # 0-100
            "motivo": str,     # Explicacao breve
            "pode_resolver": bool,  # Se deve tentar resolver
            "status_skip": str | None  # Status para log se pular
        }
    """
    client = anthropic.Anthropic(api_key=api_key)

    nome_tarefa = tarefa.get("nome", "")
    instrucoes = tarefa.get("instrucoes", "")
    texto_extraido = tarefa.get("texto_extraido", "")

    # Monta contexto completo
    contexto = f"""NOME DA TAREFA: {nome_tarefa}

INSTRUCOES: {instrucoes}"""

    if texto_extraido:
        # Limita texto extraido para nao estourar tokens
        texto_limitado = texto_extraido[:3000]
        contexto += f"""

CONTEUDO DOS ANEXOS (resumo):
{texto_limitado}"""

    prompt = f"""Analise esta tarefa educacional e classifique em UMA das categorias:

{contexto}

CATEGORIAS:

1. RESOLVIVEL - Tarefa que pode ser respondida E ENVIADA diretamente com texto, codigo ou documento
   Exemplos: dissertacao, exercicios, programacao, pesquisa, analise de caso, relatorio
   A entrega e feita anexando arquivo ou texto na propria plataforma

2. RESOLVIVEL_MANUAL - Tarefa que PODEMOS RESOLVER (gerar o codigo/arquivo), mas o ENVIO precisa ser manual
   Exemplos: "crie o codigo e envie o link do repositorio GitHub", "desenvolva o projeto e compartilhe no GitHub"
   Usamos quando: podemos gerar o arquivo, mas a entrega exige repositorio/link externo
   O sistema vai GERAR o arquivo e ANEXAR, mas NAO vai clicar em entregar

3. RESOLVER_PARCIAL - Tarefa que pede exercicios + FOTOS/PRINTS do aluno
   Exemplos: "resolva os exercicios e tire fotos", "faca e envie print/foto do resultado"
   Usamos quando: podemos RESOLVER os exercicios, mas o aluno precisa adicionar fotos/prints proprios
   O sistema vai RESOLVER os exercicios, gerar PDF com resolucoes, ANEXAR, mas NAO enviar
   O aluno depois adiciona suas fotos e envia manualmente
   IMPORTANTE: Se pede apenas "salve em PDF" sem mencionar fotos, e RESOLVIVEL normal

4. CERTIFICADO - Exige documento pessoal do aluno que ele precisa ter
   Exemplos: upload de certificado de curso, comprovante de atividade extracurricular, declaracao pessoal

5. AVISO - Apenas comunicado informativo, NAO requer entrega
   Exemplos: lembrete de prova, informativo sobre aula, orientacoes gerais, aviso de ferias

6. GRUPO - Requer formacao de equipe ou decisao coletiva
   Exemplos: escolher grupo, definir tema com a equipe, cadastro de integrantes

7. RECURSO_EXTERNO - Requer algo que o ALUNO JA PRECISA TER pronto OU ferramenta visual externa
   Exemplos: "envie o link do SEU portfolio existente", "link do SEU video JA GRAVADO no YouTube"
   TAMBEM INCLUI tarefas que exigem FERRAMENTAS VISUAIS/DESIGN que nao conseguimos gerar:
   - Wireframes/rabiscoframes/mockups usando: Excalidraw, Balsamiq, Figma, Adobe XD, Sketch, Canva
   - Diagramas interativos: Draw.io, Lucidchart, Miro, Whimsical
   - Prototipos de interface, telas de aplicativo, fluxogramas visuais
   - Qualquer tarefa que peca "entre 5 e 10 telas", "crie wireframes", "modele a interface"
   IMPORTANTE: Se podemos CRIAR o conteudo mas precisa de repo, use RESOLVIVEL_MANUAL
   IMPORTANTE: NAO e recurso externo se for apenas acessar material do professor

8. PRESENCIAL - Requer presenca fisica ou acao impossivel remotamente
   Exemplos: prova presencial, visita tecnica, apresentacao ao vivo

9. PESSOAL - Requer experiencia ou opiniao UNICA e PESSOAL do aluno
   Exemplos: "descreva SUA experiencia de estagio", "conte sobre SEU projeto pessoal", autoavaliacao
   IMPORTANTE: Perguntas genericas de opiniao (ex: "o que voce acha sobre X") SAO resolviveis

10. INCERTO - Instrucoes confusas, incompletas ou ambiguas demais

REGRAS DE DECISAO:
- Se parece uma tarefa academica normal (exercicio, trabalho, prova), e RESOLVIVEL
- Se pede QUALQUER tipo de documento/texto/codigo como resposta E pode anexar direto, e RESOLVIVEL
- SE PEDE CRIAR CODIGO/ARQUIVO + ENTREGAR VIA REPOSITORIO/GITHUB, e RESOLVIVEL_MANUAL
- Se pede exercicios + "tire fotos", "envie prints", "fotografe", "capture a tela" = RESOLVER_PARCIAL
- Se pede algo que o aluno JA PRECISA TER (video gravado, portfolio existente), e RECURSO_EXTERNO
- Se menciona ferramentas como Excalidraw, Balsamiq, Figma, Canva, Draw.io, Miro para criar wireframes/mockups/prototipos, e RECURSO_EXTERNO
- Se pede "crie X telas", "modele a interface", "faca wireframes/rabiscoframes", e RECURSO_EXTERNO
- Apenas classifique como nao-resolvivel se tiver CERTEZA que se encaixa nas outras categorias
- Na duvida entre RESOLVIVEL e RESOLVIVEL_MANUAL, escolha RESOLVIVEL_MANUAL se mencionar repositorio

Responda APENAS com JSON valido, sem markdown:
{{"categoria": "CATEGORIA", "confianca": 0-100, "motivo": "explicacao breve em 1 linha"}}"""

    try:
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=200,
            messages=[{"role": "user", "content": prompt}]
        )

        resposta_texto = response.content[0].text.strip()

        # Tenta extrair JSON
        # Remove possivel markdown
        resposta_texto = re.sub(r'^```json\s*', '', resposta_texto)
        resposta_texto = re.sub(r'\s*```$', '', resposta_texto)

        resultado = json.loads(resposta_texto)

        categoria = resultado.get("categoria", CATEGORIA_INCERTO).upper()
        confianca = int(resultado.get("confianca", 50))
        motivo = resultado.get("motivo", "")

        # Valida categoria
        categorias_validas = [
            CATEGORIA_RESOLVIVEL, CATEGORIA_RESOLVIVEL_MANUAL, CATEGORIA_RESOLVER_PARCIAL,
            CATEGORIA_CERTIFICADO, CATEGORIA_AVISO,
            CATEGORIA_GRUPO, CATEGORIA_RECURSO_EXTERNO, CATEGORIA_PRESENCIAL,
            CATEGORIA_PESSOAL, CATEGORIA_INCERTO
        ]
        if categoria not in categorias_validas:
            categoria = CATEGORIA_INCERTO
            confianca = 30

        # Determina se pode resolver
        pode_resolver = True
        status_skip = None
        anexar_apenas = False

        # Categorias que pulamos completamente
        if categoria in CATEGORIAS_PULAR and confianca >= 70:
            pode_resolver = False
            status_map = {
                CATEGORIA_CERTIFICADO: "skipped_certificate",
                CATEGORIA_AVISO: "skipped_announcement",
                CATEGORIA_GRUPO: "skipped_group",
                CATEGORIA_RECURSO_EXTERNO: "skipped_external",
                CATEGORIA_PRESENCIAL: "skipped_presence",
                CATEGORIA_PESSOAL: "skipped_personal",
            }
            status_skip = status_map.get(categoria, "skipped")

        # Categorias que resolvemos mas não enviamos
        if categoria in CATEGORIAS_ANEXAR_APENAS:
            pode_resolver = True
            anexar_apenas = True

        # Se confianca muito baixa, marca como incerto
        if confianca < 40:
            pode_resolver = False
            status_skip = "skipped_uncertain"

        logger.info(f"Analise de intencao: {categoria} ({confianca}%) - {motivo}")

        return {
            "categoria": categoria,
            "confianca": confianca,
            "motivo": motivo,
            "pode_resolver": pode_resolver,
            "anexar_apenas": anexar_apenas,
            "status_skip": status_skip,
            "flag_revisar": categoria == CATEGORIA_INCERTO or (40 <= confianca < 70),
        }

    except json.JSONDecodeError as e:
        logger.warning(f"Erro ao parsear JSON da analise: {e}")
        # Em caso de erro, assume resolvivel para nao bloquear
        return {
            "categoria": CATEGORIA_RESOLVIVEL,
            "confianca": 50,
            "motivo": "Erro na analise, assumindo resolvivel",
            "pode_resolver": True,
            "anexar_apenas": False,
            "status_skip": None,
            "flag_revisar": True,
        }

    except Exception as e:
        logger.error(f"Erro na analise de intencao: {e}")
        # Em caso de erro, assume resolvivel
        return {
            "categoria": CATEGORIA_RESOLVIVEL,
            "confianca": 50,
            "motivo": f"Erro: {str(e)[:50]}",
            "pode_resolver": True,
            "anexar_apenas": False,
            "status_skip": None,
            "flag_revisar": True,
        }


# ============================================================
# CLAUDE API
# ============================================================

def resolver_com_claude(tarefa: dict, api_key: str, nome_aluno: str = "") -> str | None:
    """Resolve tarefa usando Claude API com retry."""
    client = anthropic.Anthropic(api_key=api_key)

    formatos_str = ", ".join(TODOS_FORMATOS)

    # Instrucao sobre nome do aluno
    nome_instrucao = ""
    if nome_aluno:
        nome_instrucao = f"\nSeu nome e: {nome_aluno}. Use este nome se a tarefa exigir identificacao."

    # Monta bloco de texto extraido dos anexos (se houver)
    texto_extraido = tarefa.get("texto_extraido", "").strip()
    bloco_texto_extraido = ""
    if texto_extraido:
        bloco_texto_extraido = f"""

CONTEUDO EXTRAIDO DOS ANEXOS:
{texto_extraido}
"""

    prompt = f"""Voce e um estudante universitario resolvendo suas proprias tarefas.
Escreva de forma NATURAL e HUMANA, como um aluno real escreveria.{nome_instrucao}

TAREFA: {tarefa.get('nome', 'Sem nome')}

INSTRUCOES DA TAREFA: {tarefa.get('instrucoes', 'Nao especificadas')}
{bloco_texto_extraido}
Analise o conteudo dos anexos (texto extraido e/ou imagens) e resolva a tarefa.

REGRA OBRIGATORIA - PRIMEIRA LINHA:
A primeira linha da sua resposta DEVE ser EXATAMENTE no formato:
[FORMATO: xxx]
Onde xxx e um dos seguintes: {formatos_str}
Escolha o formato baseado no que a tarefa pede como entrega.
Exemplos:
- Se pede pagina(s) HTML -> [FORMATO: html]
- Se pede documento/trabalho/pesquisa/relatorio Word -> [FORMATO: docx]
- Se pede planilha/tabela Excel -> [FORMATO: xlsx]
- Se pede apresentacao/slides PowerPoint -> [FORMATO: pptx]
- Se pede codigo Python -> [FORMATO: py]
- Se pede codigo JavaScript -> [FORMATO: js]
- Se pede codigo TypeScript -> [FORMATO: ts]
- Se pede codigo Java (nao Android) -> [FORMATO: java]
- Se pede codigo Kotlin -> [FORMATO: kotlin]
- Se pede codigo C -> [FORMATO: c]
- Se pede codigo C++ -> [FORMATO: cpp]
- Se pede CSS/folha de estilo -> [FORMATO: css]
- Se pede SQL/consultas banco de dados -> [FORMATO: sql]
- Se pede multiplos arquivos (ex: HTML+CSS+JS) -> [FORMATO: zip]
- Se pede app/projeto Android Studio -> [FORMATO: android]
- Se pede resposta em texto simples na caixa -> [FORMATO: texto]

REGRAS DA RESPOSTA:
1. Escreva como um ESTUDANTE REAL, nao como uma IA - use linguagem natural e informal quando apropriado
2. NAO use markdown (###, **, etc) em documentos DOCX - escreva texto corrido normal
3. Para DOCX: use paragrafos normais, nao use simbolos especiais
4. Para XLSX: formate como tabela usando | para colunas
5. Para PPTX: separe slides com --- e use texto normal
6. Para CODIGO: escreva codigo LIMPO e FUNCIONAL, com POUCOS comentarios (apenas onde realmente necessario)
7. NAO inclua explicacoes sobre o que voce fez, apenas entregue o conteudo solicitado
8. NAO seja excessivamente formal ou robotico - seja natural como um aluno
9. NUNCA use placeholders como [Seu Nome], [Nome], [Seu Curso], [Data], [Professor], etc. Va direto ao conteudo da tarefa. Se a tarefa EXIGIR nome, use o nome informado acima
10. Para ANDROID: gere cada arquivo em bloco separado:
   - Codigo Java em ```java ... ``` ou Kotlin em ```kotlin ... ```
   - Layouts XML em ```xml ... ``` (activity_main.xml, etc)
   - NAO inclua AndroidManifest.xml nem build.gradle (serao gerados automaticamente)
   - Inclua imports necessarios no codigo Java/Kotlin"""

    content = [{"type": "text", "text": prompt}]

    for screenshot_path in tarefa.get("screenshots", [])[:10]:
        try:
            with open(screenshot_path, "rb") as f:
                image_data = base64.standard_b64encode(f.read()).decode("utf-8")
            content.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/png",
                    "data": image_data
                }
            })
        except Exception:
            pass

    # Adiciona PDFs nativos (suporte direto da API Claude)
    for pdf_b64 in tarefa.get("pdf_base64", []):
        try:
            content.append({
                "type": "document",
                "source": {
                    "type": "base64",
                    "media_type": "application/pdf",
                    "data": pdf_b64
                }
            })
            logger.info("PDF nativo adicionado ao request do Claude")
        except Exception as e:
            logger.error(f"Erro ao adicionar PDF nativo: {e}")

    # Lista de modelos para tentar (fallback)
    modelos = [
        ("claude-sonnet-4-20250514", "Sonnet 4"),
        ("claude-haiku-4-5-20251001", "Haiku 4.5"),
    ]

    for modelo_id, modelo_nome in modelos:
        logger.info(f"Tentando com {modelo_nome}...")

        for tentativa in range(1, CLAUDE_MAX_RETRIES + 1):
            try:
                response = client.messages.create(
                    model=modelo_id,
                    max_tokens=8096,
                    messages=[{"role": "user", "content": content}]
                )
                logger.info(f"Sucesso com {modelo_nome}!")
                return response.content[0].text
            except Exception as e:
                logger.error(f"Erro na API {modelo_nome} (tentativa {tentativa}/{CLAUDE_MAX_RETRIES}): {e}")
                if tentativa < CLAUDE_MAX_RETRIES:
                    time.sleep(tentativa * 5)

        # Se chegou aqui, falhou todas as tentativas com este modelo
        logger.warning(f"{modelo_nome} falhou, tentando proximo modelo...")

    # Todos os modelos falharam
    logger.error("Todos os modelos falharam!")
    return None

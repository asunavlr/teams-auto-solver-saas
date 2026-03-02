"""
Modulo de extracao de texto de arquivos baixados.
Suporta PDF, DOCX, XLSX e PPTX.
"""

import base64
from pathlib import Path
from loguru import logger


def extrair_conteudo_pdf(filepath: Path) -> dict:
    """
    Extrai texto e base64 de um PDF.

    Returns:
        {
            "texto": str com texto extraido,
            "base64_data": str com PDF em base64 (para envio nativo ao Claude),
            "paginas": int numero de paginas
        }
    """
    import pdfplumber

    texto_paginas = []
    num_paginas = 0

    try:
        with pdfplumber.open(filepath) as pdf:
            num_paginas = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                texto_pagina = page.extract_text() or ""

                # Extrai tabelas separadamente para melhor formatacao
                tabelas = page.extract_tables()
                texto_tabelas = ""
                for tabela in tabelas:
                    linhas_formatadas = []
                    for linha in tabela:
                        celulas = [str(c).strip() if c else "" for c in linha]
                        linhas_formatadas.append(" | ".join(celulas))
                    texto_tabelas += "\n".join(linhas_formatadas) + "\n"

                # Combina texto e tabelas
                conteudo_pagina = texto_pagina
                if texto_tabelas and texto_tabelas.strip() not in texto_pagina:
                    conteudo_pagina += f"\n\n[TABELA]\n{texto_tabelas}"

                if conteudo_pagina.strip():
                    texto_paginas.append(f"--- Pagina {i+1} ---\n{conteudo_pagina}")

    except Exception as e:
        logger.error(f"Erro ao extrair texto do PDF: {e}")

    # Gera base64 do PDF inteiro
    base64_data = None
    try:
        tamanho = filepath.stat().st_size
        if tamanho <= 20 * 1024 * 1024:  # Limite de 20MB
            with open(filepath, "rb") as f:
                base64_data = base64.standard_b64encode(f.read()).decode("utf-8")
        else:
            logger.warning(f"PDF muito grande ({tamanho / 1024 / 1024:.1f}MB), pulando base64")
    except Exception as e:
        logger.error(f"Erro ao gerar base64 do PDF: {e}")

    return {
        "texto": "\n\n".join(texto_paginas),
        "base64_data": base64_data,
        "paginas": num_paginas,
    }


def extrair_conteudo_docx(filepath: Path) -> dict:
    """
    Extrai texto de um arquivo DOCX.

    Returns:
        {"texto": str, "paginas": int}
    """
    from docx import Document

    texto_partes = []

    try:
        doc = Document(str(filepath))

        # Extrai paragrafos
        for para in doc.paragraphs:
            if para.text.strip():
                texto_partes.append(para.text)

        # Extrai tabelas
        for tabela in doc.tables:
            linhas_formatadas = []
            for linha in tabela.rows:
                celulas = [cell.text.strip() for cell in linha.cells]
                linhas_formatadas.append(" | ".join(celulas))
            if linhas_formatadas:
                texto_partes.append("\n[TABELA]\n" + "\n".join(linhas_formatadas))

    except Exception as e:
        logger.error(f"Erro ao extrair texto do DOCX: {e}")

    return {
        "texto": "\n".join(texto_partes),
        "paginas": max(1, len(texto_partes) // 30),  # Estimativa
    }


def extrair_conteudo_xlsx(filepath: Path) -> dict:
    """
    Extrai dados de um arquivo XLSX.

    Returns:
        {"texto": str, "paginas": int}
    """
    from openpyxl import load_workbook

    texto_partes = []

    try:
        wb = load_workbook(str(filepath), data_only=True)

        for nome_planilha in wb.sheetnames:
            ws = wb[nome_planilha]
            texto_partes.append(f"=== Planilha: {nome_planilha} ===")

            for row in ws.iter_rows(values_only=True):
                celulas = [str(c).strip() if c is not None else "" for c in row]
                if any(celulas):  # Pula linhas totalmente vazias
                    texto_partes.append(" | ".join(celulas))

        wb.close()
    except Exception as e:
        logger.error(f"Erro ao extrair dados do XLSX: {e}")

    return {
        "texto": "\n".join(texto_partes),
        "paginas": len(texto_partes) // 30 or 1,
    }


def extrair_conteudo_pptx(filepath: Path) -> dict:
    """
    Extrai texto de um arquivo PPTX.

    Returns:
        {"texto": str, "paginas": int}
    """
    from pptx import Presentation

    texto_partes = []
    num_slides = 0

    try:
        prs = Presentation(str(filepath))
        num_slides = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            textos_slide = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            textos_slide.append(para.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        celulas = [cell.text.strip() for cell in row.cells]
                        textos_slide.append(" | ".join(celulas))

            if textos_slide:
                texto_partes.append(f"--- Slide {i+1} ---\n" + "\n".join(textos_slide))

    except Exception as e:
        logger.error(f"Erro ao extrair texto do PPTX: {e}")

    return {
        "texto": "\n\n".join(texto_partes),
        "paginas": num_slides,
    }


def extrair_conteudo_arquivo(filepath: Path) -> dict | None:
    """
    Dispatcher: extrai conteudo baseado na extensao do arquivo.

    Returns:
        dict com {texto, paginas, base64_data(opcional)} ou None se formato nao suportado
    """
    ext = filepath.suffix.lower()

    extratores = {
        ".pdf": extrair_conteudo_pdf,
        ".docx": extrair_conteudo_docx,
        ".doc": extrair_conteudo_docx,
        ".xlsx": extrair_conteudo_xlsx,
        ".xls": extrair_conteudo_xlsx,
        ".pptx": extrair_conteudo_pptx,
        ".ppt": extrair_conteudo_pptx,
    }

    extrator = extratores.get(ext)
    if not extrator:
        logger.warning(f"Formato nao suportado para extracao: {ext}")
        return None

    logger.info(f"Extraindo conteudo de {filepath.name} ({ext})")
    resultado = extrator(filepath)
    logger.info(f"Extraido: {len(resultado.get('texto', ''))} caracteres, {resultado.get('paginas', 0)} paginas")
    return resultado
